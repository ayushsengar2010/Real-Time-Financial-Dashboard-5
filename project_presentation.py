from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Real-Time Financial Insights Dashboard"
subtitle = slide.placeholders[1]
subtitle.text = "A Modern Web App for Stock Market Analysis\nYour Name | Date"

# Slide 2: Project Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Project Overview"
content = (
    "This project is a full-stack web application designed to provide real-time insights into the stock market. "
    "It features a modern React frontend and a FastAPI backend, offering users live market data, AI-powered financial analysis, portfolio management, and an interactive chatbot for stock and trading questions."
)
slide.placeholders[1].text = content

# Slide 3: Key Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Key Features"
content = (
    "- Real-time stock data streaming from Yahoo Finance\n"
    "- AI-powered financial analysis and recommendations\n"
    "- Portfolio management and tracking\n"
    "- Alert notifications\n"
    "- User authentication (JWT & OAuth)\n"
    "- Interactive chatbot for stock/trading questions"
)
slide.placeholders[1].text = content

# Slide 4: Technology Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Technology Stack"
content = (
    "Frontend: React, TypeScript, Material-UI, Recharts\n"
    "Backend: FastAPI, SQLAlchemy, SQLite, yfinance\n"
    "AI: Mocked GPT-4 style analysis\n"
    "Deployment: Vercel (frontend), Docker-ready backend"
)
slide.placeholders[1].text = content

# Slide 5: User Authentication
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "User Authentication"
content = (
    "The app provides secure login and registration using JWT tokens. OAuth integration allows users to sign in with Google or GitHub. "
    "Logout functionality ensures users are redirected to the login page with a confirmation message. All dashboard routes are protected."
)
slide.placeholders[1].text = content

# Slide 6: Dashboard UI
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Dashboard UI"
content = (
    "The dashboard features a modern, responsive design with a sidebar for user info and quick stats. "
    "Animated cards display live market data, and interactive charts provide a clear overview of market trends and portfolio performance."
)
slide.placeholders[1].text = content
slide.shapes.add_picture('screenshots/dashboard.png', Inches(5), Inches(2), width=Inches(3))

# Slide 7: AI Financial Assistant
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "AI Financial Assistant"
content = (
    "Users can ask questions about stocks, trading, or markets. The AI provides analysis, recommendations, and risk assessment, strictly answering only finance-related questions."
)
slide.placeholders[1].text = content
slide.shapes.add_picture('screenshots/ai_assistant.png', Inches(5), Inches(2), width=Inches(3))

# Slide 8: Floating Chatbot
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Floating Chatbot"
content = (
    "A floating chatbot is always accessible on the dashboard. It answers only stock, trading, or finance questions and opens Yahoo Finance for any stock card clicked."
)
slide.placeholders[1].text = content
slide.shapes.add_picture('screenshots/chatbot.png', Inches(5), Inches(2), width=Inches(3))

# Slide 9: Deployment & Usage
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Deployment & Usage"
content = (
    "The frontend is deployed on Vercel, while the backend can be deployed on Render, Railway, or a VPS. "
    "Environment variables are used for API URLs. The app is easy to run locally with Docker or npm scripts."
)
slide.placeholders[1].text = content

# Slide 10: Demo & Screenshots
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Demo & Screenshots"
content = (
    "Below are screenshots of the Home, Login, Dashboard, and Chatbot.\n"
    "Live demo: [your Vercel URL]\nThank you! Questions?"
)
slide.placeholders[1].text = content
slide.shapes.add_picture('screenshots/demo.png', Inches(5), Inches(2), width=Inches(3))

prs.save('RealTime_Financial_Dashboard_Presentation.pptx') 